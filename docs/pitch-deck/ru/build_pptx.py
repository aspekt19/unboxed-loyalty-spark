#!/usr/bin/env python3
"""Генератор инвестиционной презентации Loyal Spark (RU) -> .pptx

Запуск:  python3 build_pptx.py
Результат: LoyalSpark_Pitch_RU.pptx (16:9) в этой же папке.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Палитра (в тон HTML-деку) ----
BG      = RGBColor(0x07, 0x0B, 0x18)
BG2     = RGBColor(0x0C, 0x14, 0x30)
PANEL   = RGBColor(0x13, 0x1C, 0x3A)
PANEL2  = RGBColor(0x18, 0x23, 0x47)
TEXT    = RGBColor(0xEE, 0xF2, 0xFF)
MUTED   = RGBColor(0x9A, 0xA7, 0xC7)
BRAND   = RGBColor(0x5B, 0x8C, 0xFF)
BRAND2  = RGBColor(0x8A, 0x6B, 0xFF)
ACCENT  = RGBColor(0x27, 0xE0, 0xB3)
GOLD    = RGBColor(0xFF, 0xCC, 0x66)
DANGER  = RGBColor(0xFF, 0x7A, 0x90)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0x2A, 0x35, 0x5E)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def solid(shape, color, line_color=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    shape.shadow.inherit = False


def no_fill(shape):
    shape.fill.background()
    shape.line.fill.background()
    shape.shadow.inherit = False


def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    solid(r, color)
    r._element.addprevious(r._element)  # keep at back-ish
    return r


def grad_bg(slide):
    """Тёмный градиентный фон."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.line.fill.background()
    r.shadow.inherit = False
    fill = r.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = BG2
    stops[1].position = 1.0
    stops[1].color.rgb = BG
    try:
        fill.gradient_angle = 60.0
    except Exception:
        pass
    return r


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; каждый параграф = list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Inter"
    return tb


def kicker(slide, text):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.62),
                                  Inches(0.018 * len(text) + 0.55), Inches(0.42))
    solid(pill, PANEL, line_color=BRAND, line_w=Pt(0.75))
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BRAND; r.font.name = "Inter"


def title(slide, text, y=Inches(1.2), size=34, color=TEXT):
    textbox(slide, Inches(0.9), y, Inches(11.5), Inches(1.0),
            [[(text, size, color, True)]], line_spacing=1.0)


def subtitle(slide, text, y=Inches(2.0)):
    textbox(slide, Inches(0.9), y, Inches(11.0), Inches(0.7),
            [[(text, 15, MUTED, False)]], line_spacing=1.1)


def footer(slide, page):
    textbox(slide, Inches(0.9), Inches(7.05), Inches(4), Inches(0.35),
            [[("Loyal", 11, TEXT, True), ("Spark", 11, BRAND, True),
              ("  ·  Onchain Loyalty Protocol", 11, MUTED, False)]])
    textbox(slide, Inches(10.9), Inches(7.05), Inches(1.6), Inches(0.35),
            [[(f"{page:02d} / 17", 11, MUTED, False)]], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=PANEL, line=LINE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    solid(c, fill, line_color=line, line_w=Pt(1))
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    return c


def card_text(slide, x, y, w, h, icon, head, body, head_color=TEXT):
    card(slide, x, y, w, h)
    pad = Inches(0.28)
    paras = []
    if icon:
        paras.append([(icon, 22, BRAND, False)])
    paras.append([(head, 16, head_color, True)])
    paras.append([(body, 12.5, MUTED, False)])
    textbox(slide, x + pad, y + Inches(0.22), w - pad * 2, h - Inches(0.4),
            paras, line_spacing=1.12, space_after=6)


# ============ SLIDE 1 — COVER ============
s = prs.slides.add_slide(BLANK)
grad_bg(s)
# decorative accent bars
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(0.12))
solid(bar, BRAND)
textbox(s, Inches(0.9), Inches(0.95), Inches(8), Inches(0.5),
        [[("⚡ ONCHAIN LOYALTY PROTOCOL · BASE MAINNET", 13, BRAND, True)]])
textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6),
        [[("Loyal Spark", 76, WHITE, True)]])
textbox(s, Inches(0.9), Inches(3.6), Inches(10.5), Inches(1.2),
        [[("Программы лояльности нового поколения для бизнеса —", 24, TEXT, False)],
         [("управляемые людьми и AI-агентами.", 24, TEXT, True)]], line_spacing=1.15)
# badges
badges = ["🔗 Построено на Base", "💵 USDC-native", "🤖 REST + MCP для агентов", "🌐 loyalspark.online"]
bx = Inches(0.9)
for b in badges:
    w = Inches(0.085 * len(b) + 0.5)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.4), w, Inches(0.5))
    solid(pill, PANEL, line_color=LINE, line_w=Pt(0.75))
    try: pill.adjustments[0] = 0.5
    except Exception: pass
    tf = pill.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = b; r.font.size = Pt(12); r.font.color.rgb = MUTED; r.font.name = "Inter"
    bx = bx + w + Inches(0.18)
textbox(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
        [[("Презентация для инвесторов и корпоративных партнёров · 2026", 12, MUTED, False)]])

# ============ SLIDE 2 — PROBLEM ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "01 · Проблема")
title(s, "Лояльность застряла в прошлом веке")
subtitle(s, "Традиционные программы дороги для бизнеса, бесполезны для клиента и закрыты для автоматизации.")
cw, gap = Inches(3.7), Inches(0.25)
cx = Inches(0.9); cy = Inches(2.9); ch = Inches(3.0)
card_text(s, cx, cy, cw, ch, "💸", "Дорого для бизнеса",
          "Legacy-SaaS берёт $99–499/мес или удерживает 15–30% с оборота программы.")
card_text(s, cx + cw + gap, cy, cw, ch, "🔒", "Баллы не принадлежат клиенту",
          "Это строка в чужой базе данных: их обесценивают, «сжигают» по сроку и нельзя перенести.")
card_text(s, cx + (cw + gap) * 2, cy, cw, ch, "🤖", "Закрыто для агентов",
          "Программы изолированы, нет API для AI-агентов и машинных платежей — главного тренда 2025–2027.")
footer(s, 2)

# ============ SLIDE 3 — SOLUTION ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "02 · Решение")
title(s, "Двухрежимная платформа лояльности")
subtitle(s, "Один протокол — два интерфейса: люди через портал, AI-агенты через API и MCP. Общие смарт-контракты, токены и БД.")
card_text(s, cx, cy, cw, ch, "🪙", "Владейте программой",
          "Разверните свой ERC-20 токен лояльности на Base. Эмиссия, переводы и погашение — onchain.")
card_text(s, cx + cw + gap, cy, cw, ch, "💵", "Платите за использование",
          "Предсказуемые тарифы в USDC. Без процента с оборота и скрытых блокировок.")
card_text(s, cx + (cw + gap) * 2, cy, cw, ch, "⚙️", "Автоматизируйте агентами",
          "AI-агенты ведут программы через REST + MCP, оплачивая вызовы микроплатежами x402 / MPP.")
footer(s, 3)

# ============ SLIDE 4 — WHY NOW ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "03 · Почему сейчас")
title(s, "Сходятся три волны")
items = [
    ("🤖", "Экономика AI-агентов.", "Агенты уже совершают покупки и платят за API. Нужны программируемые сервисы с машинной оплатой — мы такими родились."),
    ("⛓️", "Стейблкоины и L2 стали зрелыми.", "Base + USDC дают мгновенные расчёты с копеечной комиссией — onchain-лояльность наконец рентабельна."),
    ("📉", "Давление на маржу ритейла.", "Бизнес ищет дешёвые инструменты удержания. Мы убираем процент с оборота и дорогую SaaS-абонплату."),
]
yy = Inches(2.3)
for ic, head, body in items:
    c = card(s, Inches(0.9), yy, Inches(11.5), Inches(1.25))
    textbox(s, Inches(1.15), yy + Inches(0.28), Inches(0.7), Inches(0.7), [[(ic, 26, BRAND, False)]])
    textbox(s, Inches(2.0), yy + Inches(0.22), Inches(10.1), Inches(0.95),
            [[(head + " ", 17, WHITE, True)], [(body, 13.5, MUTED, False)]], line_spacing=1.1)
    yy = yy + Inches(1.45)
footer(s, 4)

# ============ SLIDE 5 — MARKET ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "04 · Рынок")
title(s, "Большой и недоинвестированный")
stats = [
    ("$200B+", "TAM — глобальный рынок лояльности (SaaS + обязательства по баллам)"),
    ("$60B", "SAM — цифровая лояльность SMB и mid-market"),
    ("630K+", "SMB с программами лояльности только в США"),
    ("$100B+", "баллов выпускается и не погашается ежегодно"),
]
sw = Inches(2.78); sx = Inches(0.9); sy = Inches(2.7)
for i, (num, lab) in enumerate(stats):
    x = sx + i * (sw + Inches(0.16))
    card(s, x, sy, sw, Inches(2.7))
    textbox(s, x + Inches(0.2), sy + Inches(0.4), sw - Inches(0.4), Inches(0.9),
            [[(num, 38, BRAND, True)]], align=PP_ALIGN.LEFT)
    textbox(s, x + Inches(0.2), sy + Inches(1.4), sw - Inches(0.4), Inches(1.1),
            [[(lab, 13, MUTED, False)]], line_spacing=1.12)
textbox(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.6),
        [[("Источники: отраслевые отчёты по рынку лояльности и SMB-коммерции. SAM — обслуживаемый сегмент цифровой лояльности.", 11.5, MUTED, False)]])
footer(s, 5)

# ============ SLIDE 6 — PRODUCT ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "05 · Продукт")
title(s, "Два портала + один программируемый API")
def flow_card(x, title_txt, steps):
    c = card(s, x, Inches(2.55), Inches(5.6), Inches(4.0))
    textbox(s, x + Inches(0.3), Inches(2.75), Inches(5.0), Inches(0.5),
            [[(title_txt, 17, WHITE, True)]])
    yy = Inches(3.35)
    for n, (h, b) in enumerate(steps, 1):
        num = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), yy, Inches(0.42), Inches(0.42))
        solid(num, BRAND)
        try: num.adjustments[0] = 0.25
        except Exception: pass
        tfn = num.text_frame; tfn.margin_top=0; tfn.margin_bottom=0
        pn = tfn.paragraphs[0]; pn.alignment=PP_ALIGN.CENTER
        rn = pn.add_run(); rn.text=str(n); rn.font.size=Pt(13); rn.font.bold=True; rn.font.color.rgb=WHITE
        textbox(s, x + Inches(0.9), yy - Inches(0.04), Inches(4.5), Inches(1.0),
                [[(h + ".  ", 14, TEXT, True), (b, 12.5, MUTED, False)]], line_spacing=1.08)
        yy = yy + Inches(1.0)
flow_card(Inches(0.9), "🏪 Сценарий бизнеса", [
    ("Запуск", "ERC-20 токен лояльности на Base за минуты"),
    ("Работа", "Начисление по email/телефону/кошельку, ваучеры, уровни, RFM"),
    ("Команда и агенты", "Филиалы, сотрудники, AI-агенты со scoped API-ключами"),
])
flow_card(Inches(6.85), "👤 Сценарий клиента", [
    ("Получение", "Токены в смарт-кошелёк (Privy / SIWE) — без знаний о крипте"),
    ("Погашение", "Тратит на награды или обменивает P2P через onchain-эскроу"),
    ("Владение", "Токены в кошельке клиента, полностью переносимы"),
])
footer(s, 6)

# ============ SLIDE 7 — TECH / SECURITY ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "06 · Технологии и безопасность")
title(s, "Корпоративный уровень с первого дня")
subtitle(s, "Архитектура, которую не стыдно показать службе безопасности крупного партнёра.")
tech = [
    ("🔐", "Гибкая аутентификация", "Email / телефон / Google через Privy + встроенные кошельки; SIWE для крипто-native."),
    ("🛡️", "Защита данных", "Row Level Security на каждой таблице, scoped API-ключи lsk_/rwk_ с SHA-256."),
    ("🔑", "MPC-кошельки", "Coinbase CDP: приватные ключи не покидают защищённый анклав."),
    ("📜", "Открытые стандарты", "OpenAPI 3.1, agent.json, MCP-каталоги (Glama, Smithery, OpenServ)."),
    ("⚡", "Платежи HTTP 402", "Нативный x402 и MPP — машинная оплата без API-ключей."),
    ("🧾", "Полный аудит", "Журналирование операций, rate-limiting, Builder Code (ERC-8021)."),
]
gw, gh = Inches(3.7), Inches(1.85)
gx0, gy0 = Inches(0.9), Inches(2.85)
for i, (ic, h, b) in enumerate(tech):
    col = i % 3; row = i // 3
    x = gx0 + col * (gw + Inches(0.25)); y = gy0 + row * (gh + Inches(0.2))
    card_text(s, x, y, gw, gh, ic, h, b)
footer(s, 7)

# ============ SLIDE 8 — BUSINESS MODEL ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "07 · Бизнес-модель")
title(s, "Две оси выручки + комиссия с оборота")

def price_card(x, title_txt, rows):
    c = card(s, x, Inches(2.4), Inches(5.6), Inches(3.6))
    textbox(s, x + Inches(0.3), Inches(2.6), Inches(5.0), Inches(0.5),
            [[(title_txt, 16, WHITE, True)]])
    yy = Inches(3.25)
    for name, price, detail in rows:
        textbox(s, x + Inches(0.3), yy, Inches(1.5), Inches(0.5), [[(name, 14, TEXT, True)]])
        textbox(s, x + Inches(1.75), yy, Inches(1.3), Inches(0.5), [[(price, 14, ACCENT, True)]])
        textbox(s, x + Inches(3.0), yy, Inches(2.5), Inches(0.6), [[(detail, 11.5, MUTED, False)]], line_spacing=1.0)
        # divider
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), yy + Inches(0.62), Inches(5.0), Pt(0.75))
        solid(ln, LINE)
        yy = yy + Inches(0.82)
price_card(Inches(0.9), "🏪 Merchant SaaS (USDC на Base)", [
    ("Starter", "$39/мес", "Портал, программы, CRM-light"),
    ("Growth", "$79/мес", "Аналитика, больше мест"),
    ("Scale", "$149/мес", "Корп-бюджеты, приоритет"),
])
price_card(Inches(6.85), "🤖 AI-агенты (API + MCP)", [
    ("Free", "$0", "200 вызовов · 1.25% mint fee"),
    ("Pro", "$49/мес", "10 000 вызовов · 0.50%"),
    ("Enterprise", "$129/мес", "∞ вызовов · 0.25%"),
])
textbox(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.7),
        [[("Pay-per-call (x402 / MPP): от ~$0.001 за чтение, ~$0.005–0.05 за запись. Главный рычаг масштабирования — % mint fee с onchain-оборота.", 12, MUTED, False)]], line_spacing=1.1)
footer(s, 8)

# ============ SLIDE 9 — COMPETITION ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "08 · Конкуренция")
title(s, "Legacy-лояльность vs Loyal Spark")
# old col
oc = card(s, Inches(0.9), Inches(2.4), Inches(5.6), Inches(4.1), fill=PANEL, line=DANGER)
textbox(s, Inches(1.15), Inches(2.6), Inches(5.0), Inches(0.5), [[("Традиционные платформы", 16, DANGER, True)]])
old = ["Square Loyalty: $45–105/мес за точку", "LoyaltyLion / Yotpo: $200–700+/мес",
       "Smile.io: процент с оборота и плата за заказ", "Закрытые API, нет доступа для AI-агентов",
       "Баллы = строка в БД, нет реального владения"]
yy = Inches(3.2)
for t in old:
    textbox(s, Inches(1.15), yy, Inches(5.1), Inches(0.6), [[("✕  ", 13, DANGER, True), (t, 13, MUTED, False)]], line_spacing=1.05)
    yy = yy + Inches(0.62)
# new col
nc = card(s, Inches(6.85), Inches(2.4), Inches(5.6), Inches(4.1), fill=PANEL2, line=ACCENT)
textbox(s, Inches(7.1), Inches(2.6), Inches(5.0), Inches(0.5), [[("Loyal Spark", 16, ACCENT, True)]])
new = ["Предсказуемый USDC-SaaS от $39/мес", "ERC-20 токены — клиент реально ими владеет",
       "Готово для агентов: REST + MCP, x402 / MPP", "SIWE и Privy — без паролей и утечек email",
       "P2P-маркетплейс с эскроу между программами"]
yy = Inches(3.2)
for t in new:
    textbox(s, Inches(7.1), yy, Inches(5.1), Inches(0.6), [[("✓  ", 13, ACCENT, True), (t, 13, TEXT, False)]], line_spacing=1.05)
    yy = yy + Inches(0.62)
footer(s, 9)

# ============ SLIDE 10 — MOAT ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "09 · Защищённость")
title(s, "Почему нас сложно скопировать")
moat = [
    ("🧱", "Двухрежимность", "Люди (портал) + AI-агенты (API + MCP) в одном протоколе."),
    ("💰", "Модель «два счёта»", "SaaS бизнеса + API агентов + % mint fee — диверсификация."),
    ("⚡", "Машинные платежи", "Нативный HTTP 402 (x402/MPP) — позиция в экономике агентов."),
    ("🌐", "На Base", "Низкие комиссии, быстрый финалити, USDC-native, экосистема Coinbase."),
    ("🔓", "Открытые стандарты", "OpenAPI, agent.json, MCP-каталоги → органический приток агентов."),
    ("🪙", "Сетевой эффект LOYAL", "Hub-and-spoke ликвидность объединяет все токены лояльности."),
]
for i, (ic, h, b) in enumerate(moat):
    col = i % 3; row = i // 3
    x = gx0 + col * (gw + Inches(0.25)); y = Inches(2.4) + row * (gh + Inches(0.25))
    card_text(s, x, y, gw, gh, ic, h, b)
footer(s, 10)

# ============ SLIDE 11 — STATUS ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "10 · Статус")
title(s, "Живой продукт, не слайды")
status = [
    ("Live MVP", "Развёрнут на Base Mainnet (Chain ID 8453)"),
    ("Web · PWA · iOS/Android", "Capacitor 8 — мультиплатформа"),
    ("25 REST + 36 MCP", "инструментов для агентов (+18 recipient MCP)"),
    ("Privy + SIWE", "+ scoped API-ключи и полный аудит"),
]
for i, (num, lab) in enumerate(status):
    x = Inches(0.9) + i * (Inches(2.78) + Inches(0.16))
    card(s, x, Inches(2.5), Inches(2.78), Inches(2.5))
    textbox(s, x + Inches(0.22), Inches(2.85), Inches(2.4), Inches(1.1),
            [[(num, 19, BRAND, True)]], line_spacing=1.0)
    textbox(s, x + Inches(0.22), Inches(4.0), Inches(2.4), Inches(0.9),
            [[(lab, 12, MUTED, False)]], line_spacing=1.1)
textbox(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.9),
        [[("Смарт-контракты в продакшене: ", 13.5, MUTED, False),
          ("LoyaltyTokenFactory 0x5F3D…dA80", 13.5, TEXT, True),
          ("  ·  ", 13.5, MUTED, False),
          ("LoyalSparkERC20 0xe6BA…27C3", 13.5, TEXT, True)],
         [("Этап: pre-revenue, онбординг первых design-партнёров.", 13.5, MUTED, False)]], line_spacing=1.2)
footer(s, 11)

# ============ SLIDE 12 — GTM ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "11 · Go-to-Market")
title(s, "Сначала люди, потом канал агентов")
rows = [
    ("Сейчас", "Первые design-партнёры", "Onchain-лояльность + дашборд на Base"),
    ("Q4 2026 – Q1 2027", "10–25 платящих бизнесов", "Тарифы Starter/Growth, кейсы"),
    ("Q2–Q3 2027", "100+ бизнесов", "Канал агентов: OpenServ, MCP-каталоги, x402"),
    ("Q4 2027 – Q1 2028", "500+ бизнесов", "Мультирегиональный SaaS + выручка с агентов"),
]
# header
hy = Inches(2.5)
textbox(s, Inches(0.9), hy, Inches(2.4), Inches(0.4), [[("ЭТАП", 12, BRAND, True)]])
textbox(s, Inches(3.5), hy, Inches(3.5), Inches(0.4), [[("ЦЕЛЬ", 12, BRAND, True)]])
textbox(s, Inches(7.3), hy, Inches(5.0), Inches(0.4), [[("ФОКУС", 12, BRAND, True)]])
yy = Inches(3.05)
for phase, goal, focus in rows:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), yy - Inches(0.12), Inches(11.5), Pt(0.75)); solid(ln, LINE)
    textbox(s, Inches(0.9), yy, Inches(2.5), Inches(0.6), [[(phase, 14, ACCENT, True)]])
    textbox(s, Inches(3.5), yy, Inches(3.6), Inches(0.6), [[(goal, 14, TEXT, True)]])
    textbox(s, Inches(7.3), yy, Inches(5.0), Inches(0.6), [[(focus, 13, MUTED, False)]], line_spacing=1.0)
    yy = yy + Inches(0.85)
footer(s, 12)

# ============ SLIDE 13 — PROJECTIONS ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "12 · Прогноз")
title(s, "Экономика роста")
subtitle(s, "Сценарий на базе модели стейкинга и транзакционных комиссий.")
cols = ["Метрика", "Год 1", "Год 2", "Год 3"]
data = [
    ("Активные бизнесы", "1 000", "10 000", "50 000"),
    ("Total Value Locked", "$2M", "$25M", "$150M"),
    ("Объём транзакций", "$500K", "$10M", "$100M"),
    ("Выручка протокола", "$50K", "$1.5M", "$12M"),
]
colx = [Inches(0.9), Inches(5.0), Inches(7.7), Inches(10.4)]
hy = Inches(2.85)
for j, c in enumerate(cols):
    textbox(s, colx[j], hy, Inches(2.6), Inches(0.4), [[(c, 12, BRAND, True)]])
yy = Inches(3.4)
for r in data:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), yy - Inches(0.12), Inches(11.5), Pt(0.75)); solid(ln, LINE)
    for j, v in enumerate(r):
        bold = (j == 3) or (j == 0)
        color = TEXT if j > 0 else MUTED
        textbox(s, colx[j], yy, Inches(2.6), Inches(0.5), [[(v, 15, color, bold)]])
    yy = yy + Inches(0.72)
textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
        [[("Прогноз иллюстративный: целевая траектория при достижении PMF и активации канала агентов.", 11.5, MUTED, False)]])
footer(s, 13)

# ============ SLIDE 14 — TOKENOMICS ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "13 · Токеномика")
title(s, "LOYAL — топливо экосистемы")
# left card: allocation
card(s, Inches(0.9), Inches(2.45), Inches(5.6), Inches(4.1))
textbox(s, Inches(1.15), Inches(2.65), Inches(5.1), Inches(0.5), [[("Распределение · 10 000 000 000 LOYAL", 15, WHITE, True)]])
alloc = [("Экосистема / Treasury", "41%", 0.41), ("Token Sales", "30%", 0.30),
         ("Команда и эдвайзеры", "20%", 0.20), ("Ликвидность (IDL)", "5%", 0.05),
         ("Маркетинг / Airdrop", "4%", 0.04)]
yy = Inches(3.25)
for name, pct, frac in alloc:
    textbox(s, Inches(1.15), yy, Inches(3.0), Inches(0.35), [[(name, 12.5, TEXT, False)]])
    textbox(s, Inches(5.7), yy, Inches(0.7), Inches(0.35), [[(pct, 12.5, ACCENT, True)]], align=PP_ALIGN.RIGHT)
    track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), yy + Inches(0.32), Inches(5.2), Inches(0.14)); solid(track, LINE)
    try: track.adjustments[0]=0.5
    except Exception: pass
    fillw = Inches(5.2 * frac)
    if fillw > 0:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), yy + Inches(0.32), fillw, Inches(0.14)); solid(b, BRAND)
        try: b.adjustments[0]=0.5
        except Exception: pass
    yy = yy + Inches(0.62)
# right card: utility
card(s, Inches(6.85), Inches(2.45), Inches(5.6), Inches(4.1))
textbox(s, Inches(7.1), Inches(2.65), Inches(5.1), Inches(0.5), [[("Утилитарность и дефляция", 15, WHITE, True)]])
util = [("🔥", "Дефляция.", "8% сжигается при выводе баллов (M-token → LOYAL), защищая цену."),
        ("🏦", "Залог-стейкинг.", "Возвратный депозит для доступа и ликвидности (от $1 000)."),
        ("🗳️", "Управление DAO.", "Застейканный LOYAL даёт право голоса по параметрам протокола.")]
yy = Inches(3.3)
for ic, h, b in util:
    textbox(s, Inches(7.1), yy, Inches(0.6), Inches(0.6), [[(ic, 22, BRAND, False)]])
    textbox(s, Inches(7.75), yy, Inches(4.5), Inches(1.0), [[(h + " ", 14, WHITE, True), (b, 12.5, MUTED, False)]], line_spacing=1.1)
    yy = yy + Inches(1.05)
footer(s, 14)

# ============ SLIDE 15 — ASK ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "14 · Раунд")
textbox(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.0), [[("Seed: $3.5M – $4M", 36, WHITE, True)]])
# left: use of funds
card(s, Inches(0.9), Inches(2.45), Inches(6.6), Inches(4.1))
uf = [("Команда и таланты (14 чел.)", "70% · $2.75M", 0.70),
      ("Инфраструктура и операции", "10% · $400K", 0.10),
      ("Комьюнити", "8% · $300K", 0.08),
      ("Юридич., compliance, аудиты", "7% · $250K", 0.07),
      ("Резерв", "5% · $200K", 0.05)]
yy = Inches(2.7)
for name, val, frac in uf:
    textbox(s, Inches(1.15), yy, Inches(4.5), Inches(0.35), [[(name, 13, TEXT, False)]])
    textbox(s, Inches(5.5), yy, Inches(1.8), Inches(0.35), [[(val, 13, WHITE, True)]], align=PP_ALIGN.RIGHT)
    track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), yy + Inches(0.34), Inches(6.1), Inches(0.16)); solid(track, LINE)
    try: track.adjustments[0]=0.5
    except Exception: pass
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), yy + Inches(0.34), Inches(6.1 * frac), Inches(0.16)); solid(b, BRAND)
    try: b.adjustments[0]=0.5
    except Exception: pass
    yy = yy + Inches(0.72)
# right: valuation / runway
card(s, Inches(7.7), Inches(2.45), Inches(4.7), Inches(4.1), fill=PANEL2)
textbox(s, Inches(8.0), Inches(2.85), Inches(4.0), Inches(0.4), [[("Pre-money оценка", 13, MUTED, False)]])
textbox(s, Inches(8.0), Inches(3.25), Inches(4.0), Inches(0.9), [[("$12–15M", 40, ACCENT, True)]])
textbox(s, Inches(8.0), Inches(4.35), Inches(4.0), Inches(0.4), [[("Runway", 13, MUTED, False)]])
textbox(s, Inches(8.0), Inches(4.75), Inches(4.0), Inches(0.9), [[("18 месяцев", 36, WHITE, True)]])
textbox(s, Inches(8.0), Inches(5.85), Inches(4.0), Inches(0.6), [[("До product-market fit и готовности к Series A.", 12.5, MUTED, False)]], line_spacing=1.1)
footer(s, 15)

# ============ SLIDE 16 — MILESTONES ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
kicker(s, "15 · Цели раунда")
title(s, "Что мы достигнем за 18 месяцев")
ms = [
    ("🏪", "1 000+ бизнесов", "Онбординг SMB в нескольких вертикалях."),
    ("🔒", "$2M+ TVL", "Заблокировано в стейках бизнесов."),
    ("📈", "$500K+ оборот", "Транзакции через платформу."),
    ("🎯", "Product-Market Fit", "Подтверждён с положительной юнит-экономикой."),
    ("💹", "Путь к прибыли", "Устойчивая модель выручки."),
    ("🚀", "Series A Ready", "Сильные метрики и траектория роста."),
]
for i, (ic, h, b) in enumerate(ms):
    col = i % 3; row = i // 3
    x = gx0 + col * (gw + Inches(0.25)); y = Inches(2.4) + row * (gh + Inches(0.25))
    card_text(s, x, y, gw, gh, ic, h, b)
footer(s, 16)

# ============ SLIDE 17 — CLOSING ============
s = prs.slides.add_slide(BLANK); grad_bg(s)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.38), EMU_W, Inches(0.12)); solid(bar, BRAND)
textbox(s, Inches(0.9), Inches(1.4), Inches(8), Inches(0.5), [[("СПАСИБО", 13, BRAND, True)]])
textbox(s, Inches(0.9), Inches(2.2), Inches(11.8), Inches(1.6),
        [[("Сделаем лояльность ", 48, WHITE, True), ("программируемой", 48, BRAND, True)]], line_spacing=1.05)
textbox(s, Inches(0.9), Inches(4.1), Inches(10.5), Inches(1.0),
        [[("Откроем onchain-лояльность для следующих 50 миллионов компаний — и для AI-агентов, которые работают рядом с ними.", 18, TEXT, False)]], line_spacing=1.3)
contacts = ["🌐 loyalspark.online", "✉️ admin@loyalspark.online", "𝕏 @Loyal_Spark"]
bx = Inches(0.9)
for c in contacts:
    w = Inches(0.1 * len(c) + 0.5)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.5), w, Inches(0.5))
    solid(pill, PANEL, line_color=LINE, line_w=Pt(0.75))
    try: pill.adjustments[0]=0.5
    except Exception: pass
    tf = pill.text_frame; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r = p.add_run(); r.text=c; r.font.size=Pt(13); r.font.color.rgb=TEXT; r.font.name="Inter"
    bx = bx + w + Inches(0.2)
textbox(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
        [[("© 2026 Loyal Spark · Построено на Base · USDC-native · Agent-ready", 11.5, MUTED, False)]])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LoyalSpark_Pitch_RU.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
